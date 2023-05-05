#Script to iterate over a file structure of pptx files and scrape the first slide for the title.
#Then save a copy of the file using that title as the filename.
#For recovery of powerpoint filenames.
#Matthew Condren 2023

import os
import re
import shutil

#Method to scan the folders and create a file list.
def fileList(extension):
    path = os.getcwd() + "\\before"
    fileList = []
    print ("path: " + path + "\n")
    for root, dirs, files in os.walk(path):
        for name in files:
            if name.find(extension) != -1:
                fileList.append(os.path.join(root, name))
    return fileList
    
# Method to take in a list of files and iterate through them.    
def read_pptx(fileList):
    
    from pptx import Presentation
    fileDict = {}
    
    for file in fileList:
        f = open(f'{file}', "rb") #open the file
        deck = Presentation(f)
        text = []
        # Only read the first slide in the deck. 
        first_slide = deck.slides[0]
        
        # Look through all the shapes on the first slide.  If it does not contain text (eg. a picture), skip it.
        # As soon as the first text box is detected, record it and stop.
        for shape in first_slide.shapes:
            if not shape.has_text_frame:
                continue
            para = shape.text_frame.paragraphs[0]
            text.append(para.runs[0].text)
            break # Stop after first text box
        
        fileDict[file] = " ".join(stripForbidden(text))
        f.close() #close the file so renaming can happen without error later
    return fileDict

#Method to modify the filename and resave file elsewhere, preserving the originals.
def copyAndRename(filesAndTitles, extension):
    
    for key in filesAndTitles.keys():
        print ("\nReading File: " + key)
        newPath = (key.split("before"))
        newPath = (newPath[0] + "after\\" + extension.split(".")[1] + "\\" + (filesAndTitles[key] + extension))
        print ("Renaming to: " + newPath)
        
        #Try to write renamed file to disk
        try:
            os.renames(key, newPath)
        except Exception as error:
            print ("Failed to rename " + key)
            print (error)

#strip out forbidden characters in new filename
def stripForbidden(name):
    forbiddenChars = [">", "<", "/", ":", '"', "|", "?", "*"]
    output = name
    for forbiddenChar in forbiddenChars:
        if name[0].find(forbiddenChar):
            name[0] = name[0].replace(forbiddenChar, "")
    return name
    
outputDict = {}
outputDict.update(read_pptx(fileList(".pptx")))
copyAndRename(outputDict, ".pptx")
  